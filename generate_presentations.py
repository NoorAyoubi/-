import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_title_slide(prs, title_text, subtitle_text, sub_info_lines, is_rtl=True):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(13, 15, 23)) # Dark slate bg
    
    # Golden horizontal line accent
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.33), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(197, 168, 128) # Gold accent
    line.line.fill.background()

    # Main Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    if is_rtl:
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
        
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle_text
    p2.font.name = 'Arial'
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(197, 168, 128)
    p2.space_before = Pt(12)
    if is_rtl:
        p2.alignment = PP_ALIGN.RIGHT
    else:
        p2.alignment = PP_ALIGN.LEFT

    # Sub-info lines (Submitter details)
    txBoxInfo = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(10.33), Inches(2))
    tfInfo = txBoxInfo.text_frame
    tfInfo.word_wrap = True
    for idx, line_text in enumerate(sub_info_lines):
        p_info = tfInfo.add_paragraph() if idx > 0 else tfInfo.paragraphs[0]
        p_info.text = line_text
        p_info.font.name = 'Arial'
        p_info.font.size = Pt(14)
        p_info.font.color.rgb = RGBColor(160, 174, 192) # Muted grey
        p_info.space_before = Pt(6)
        if is_rtl:
            p_info.alignment = PP_ALIGN.RIGHT
        else:
            p_info.alignment = PP_ALIGN.LEFT

def create_content_slide(prs, part_title, slide_title, bullets, is_rtl=True):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(13, 15, 23)) # Dark bg

    # Tiny top header (Part number)
    txPart = slide.shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(11.33), Inches(0.4))
    tfPart = txPart.text_frame
    pPart = tfPart.paragraphs[0]
    pPart.text = part_title
    pPart.font.name = 'Arial'
    pPart.font.size = Pt(12)
    pPart.font.bold = True
    pPart.font.color.rgb = RGBColor(197, 168, 128) # Gold
    if is_rtl:
        pPart.alignment = PP_ALIGN.RIGHT
    else:
        pPart.alignment = PP_ALIGN.LEFT

    # Slide Title
    txTitle = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(0.8))
    tfTitle = txTitle.text_frame
    pTitle = tfTitle.paragraphs[0]
    pTitle.text = slide_title
    pTitle.font.name = 'Arial'
    pTitle.font.size = Pt(28)
    pTitle.font.bold = True
    pTitle.font.color.rgb = RGBColor(255, 255, 255) # White
    if is_rtl:
        pTitle.alignment = PP_ALIGN.RIGHT
    else:
        pTitle.alignment = PP_ALIGN.LEFT

    # Content Area
    txContent = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    tfContent = txContent.text_frame
    tfContent.word_wrap = True

    for idx, bullet in enumerate(bullets):
        p = tfContent.add_paragraph() if idx > 0 else tfContent.paragraphs[0]
        
        # Indent management
        if bullet.startswith("  * ") or bullet.startswith("    * ") or bullet.startswith("  - ") or bullet.startswith("    - "):
            clean_text = bullet.replace("  * ", "").replace("    * ", "").replace("  - ", "").replace("    - ", "")
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(160, 174, 192) # Muted text
        else:
            clean_text = bullet.replace("* ", "").replace("- ", "")
            p.level = 0
            p.font.size = Pt(19)
            p.font.bold = True if ":" in clean_text or clean_text.startswith("בתור") or clean_text.startswith("🔥") else False
            p.font.color.rgb = RGBColor(255, 255, 255) # White text

        p.text = clean_text
        p.font.name = 'Arial'
        p.space_after = Pt(10)
        if is_rtl:
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

def create_table_slide(prs, part_title, slide_title, table_headers, table_rows, is_rtl=True):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(13, 15, 23)) # Dark bg

    # Tiny top header
    txPart = slide.shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(11.33), Inches(0.4))
    tfPart = txPart.text_frame
    pPart = tfPart.paragraphs[0]
    pPart.text = part_title
    pPart.font.name = 'Arial'
    pPart.font.size = Pt(12)
    pPart.font.bold = True
    pPart.font.color.rgb = RGBColor(197, 168, 128)
    if is_rtl:
        pPart.alignment = PP_ALIGN.RIGHT
    else:
        pPart.alignment = PP_ALIGN.LEFT

    # Slide Title
    txTitle = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(0.8))
    tfTitle = txTitle.text_frame
    pTitle = tfTitle.paragraphs[0]
    pTitle.text = slide_title
    pTitle.font.name = 'Arial'
    pTitle.font.size = Pt(28)
    pTitle.font.bold = True
    pTitle.font.color.rgb = RGBColor(255, 255, 255)
    if is_rtl:
        pTitle.alignment = PP_ALIGN.RIGHT
    else:
        pTitle.alignment = PP_ALIGN.LEFT

    # Add Table
    rows = len(table_rows) + 1
    cols = len(table_headers)
    
    # Table bounds
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(11.33)
    height = Inches(4.5)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Column styling widths
    col_widths = [Inches(2.5), Inches(3.0), Inches(1.5), Inches(1.5), Inches(2.83)] if cols == 5 else [Inches(3.33)] * cols
    for idx, w in enumerate(col_widths):
        if idx < len(table.columns):
            table.columns[idx].width = w

    # Format Headers
    for col_idx, header in enumerate(table_headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(197, 168, 128) # Gold headers
        # Text styling
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(13, 15, 23) # Dark text on gold
            if is_rtl:
                paragraph.alignment = PP_ALIGN.RIGHT
            else:
                paragraph.alignment = PP_ALIGN.LEFT

    # Format Body Rows
    for row_idx, row_data in enumerate(table_rows):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_value
            cell.fill.solid()
            # Alternating row background for modern premium feel
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(22, 27, 38)
            else:
                cell.fill.fore_color.rgb = RGBColor(17, 20, 28)
            # Text styling
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                if is_rtl:
                    paragraph.alignment = PP_ALIGN.RIGHT
                else:
                    paragraph.alignment = PP_ALIGN.LEFT

# ==========================================
# HEBREW PPTX BUILDER (8 SLIDES)
# ==========================================
def build_hebrew_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Slide 1 (Title, Problem, Hypotheses)
    create_content_slide(
        prs,
        "חלק 1: הגדרת הבעיה והשערות",
        "הגדרת הבעיה והשערות היסוד (חלק 1)",
        [
            "כותרת הפרויקט: גשר אל-קודס (JLM Legal Bridge)",
            "כותרת משנה: פלטפורמה דיגיטלית חכמה לניהול הקשר בין עורך הדין ללקוח, זימון תורים והערכת סיכויי תביעה מבוססת AI",
            "🔥 הבעיה האמיתית (Pain אמיתי):",
            "  - אצל עורך הדין: עומס שיחות והודעות לא רלוונטיות | בזבוז זמן יקר על לקוחות לא מתאימים | חזרה על אותן שאלות.",
            "  - אצל הלקוח: אינו יודע אם המקרה שלו מצדיק עורך דין | שפה משפטית מורכבת | זמני המתנה ארוכים ותשובות לא ברורות | פחד וחשש משיפוטיות.",
            "🎯 קהל היעד והצרכים:",
            "  - הלקוח: צורך להבין את סיכוייו בקלות ובמהירות מרחוק ללא מבוכה (לשם כך פותח הצ'אט-בוט החכם).",
            "  - עורך הדין: צורך בכלי לניהול תיקים ותזמון תורים מובנה (CRM) לצורך ייעול הליכי הסינון והתקשורת המהירה.",
            "💡 הנחות היסוד הראשוניות (השערות ראשוניות):",
            "  - השערת ערך (Value): פדיון זכאות אנונימי ומהיר מרחוק יעלה את נכונות הלקוחות ליצירת קשר ב-40%.",
            "  - השערת שמישות (Usability): שימוש ברכיב קולי יקל על תיאור המקרה וימנע נטישת שאלונים במובייל.",
            "  - השערת קיום (Viability): מערכת ה-CRM המובנית וניהול התורים יחסכו כ-60% מזמן הסינון והתיאום המשרדי."
        ],
        is_rtl=True
    )

    # 2. Slide 2 (Part 2 - User Research & Insights)
    create_content_slide(
        prs,
        "חלק 2: סיכום מחקר משתמשים ותובנות",
        "פרסונות, JTBD וסיפורי משתמש (חלק 2)",
        [
            "👥 פרסונות המפתח (Personas):",
            "  - הלקוח עאדל (פועל בניין, 42): חרד מעלויות ומשיפוט של עורך דין, ומחפש הערכת סיכויים מיידית ועצמאית בנייד.",
            "  - עו\"ד יוסף (שותף, 35): מתוסכל מביזבוז זמן על לקוחות לא מתאימים וחזרה על אותן שאלות שוב ושוב.",
            "🔥 נקודות כאב מרכזיות (Pain Points):",
            "  - אצל הלקוח: ייאוש ואובדן תקווה בשל אורך התהליכים הבורוקרטיים והמאמץ הרב הנדרש להצלחת התביעה.",
            "  - אצל עורך הדין: בזבוז שעות עבודה יקרות בסינון שיחות ידני עמוס וחזרה על אותן שאלות בסיסיות.",
            "🎯 ניתוח Jobs-To-Be-Done (JTBD):",
            "  - משימה 1 (הלקוח): בדיקת סיכויי התביעה הרפואית שלי בדיסקרטיות ובמהירות (הצלחה: קבלת תשובה ברורה \"מצדיק/לא מצדיק\" תוך 2 דקות).",
            "  - משימה 2 (עורך הדין): סינון לקוחות וקבלת תיק מסודר מראש (הצלחה: צמצום שיחות סרק ב-60% וקבלת פרטים ממוקדים).",
            "📝 סיפורי משתמש (User Stories):",
            "  - בתור זבון חרד שאינו בטוח אם המקרה שלו מוגדר כתאונת עבודה, כאשר אני מתבייש לפנות ישירות למשרד, אני רוצה להשתמש בצ'אט בוט אנונימי ונגיש בשפתי, כדי שאקבל אישור ראשוני לזכאותי ללא שיפוט.",
            "  - בתור משתמש רגיש המשתף פרטים רפואיים, כאשר אני מסיים את הבדיקה בצ'אט, אני רוצה לבחור באופן מפורש אם לשלוח את המידע או להשמיד אותו לצמיתות, כדי שארגיש שליטה מוחלטת במידע הפרטי שלי."
        ],
        is_rtl=True
    )

    # 3. Slide 3 (Part 3 - Market Research)
    create_content_slide(
        prs,
        "חלק 3: סיכום מחקר שוק",
        "פתרונות קיימים בשוק, מגמות וניתוח SWOT (חלק 3)",
        [
            "🔍 פתרונות קיימים ופערים בשוק:",
            "  - טפסים דיגיטליים סטטיים: קשיחים, ארוכים, רק בעברית, אחוז נטישה מעל 70% בנייד ללא התאמה אישית.",
            "  - פנייה טלפונית והגעה פיזית: דורשת מאמץ וזמן נסיעה רב, חסם כניסה פסיכולוגי גבוה שגורם לייאוש מהתהליך.",
            "  - הפער בשוק: היעדר פלטפורמה אחודה המשלבת בדיקת זכאות דינמית (AI) יחד עם לוח בקרה מובנה של עורך דין לקיצור זמני תיאום וניהול תורים.",
            "📈 מגמות טכנולוגיות מובילות:",
            "  - Conversational Legal Intake: מעבר של משרדים מובילים לצ'אט-בוטים אינטראקטיביים במובייל במקום טפסים.",
            "  - AI Eligibility Checks: שימוש באלגוריתמים לקביעת היתכנות התיק וסיכויי התביעה בזמן אמת.",
            "  - Unified Practice Management: חיבור ישיר ואינטגרטיבי בין מערכות הפנייה של המשתמשים ליומן המשרד ומערכות Outlook.",
            "📊 ניתוח SWOT של המוצר:",
            "  - חוזקות (S): פתרון קול מובנה, דו-לשוניות מלאה, סינון לקוחות ללא עילה, וניהול תורים קל למשרד.",
            "  - חולשות (W): תלות בחיבור אינטרנט של המשתמש ובאמינות הנייד.",
            "  - הזדמנויות (O): פיתוח ערוץ וואטסאפ בוט (WhatsApp Integration) להקלת הגישה עוד יותר.",
            "  - איומים (T): כניסת מתחרים גדולים לשוק המקומי או התנגדות של משרדים שמרניים לטכנולוגיה."
        ],
        is_rtl=True
    )

    # 4. Slide 4 (Part 4 - Product Vision & Value Prop)
    create_content_slide(
        prs,
        "חלק 4: חזון המוצר והצעת הערך",
        "חזון, הצעת ערך ייחודית ואימות השערות (חלק 4)",
        [
            "🎯 חזון המוצר: להסיר חסמי פחד וסרבול בירוקרטי, ולהוות את הגשר הדיגיטלי הנגיש והאמין ביותר שמחבר בין נפגעי גוף לקבלת פיצוי הוגן ומהיר.",
            "💡 הצעת הערך הייחודית (Unique Value Proposition):",
            "  - הערכת זכאות משפטית ותיאום תור מיידי מול עורך דין תוך 2 דקות בלבד, בשפתך ובצורה פשוטה, ללא עלות מראש וללא כל התחייבות – עם שליטה מלאה ופרטיות מוחלטת במידע האישי שלך.",
            "⚔️ מחקר מתחרים ישירים (נקודות החוזק שלנו):",
            "  - דו-לשוניות מלאה ודינמית: (מתחרים: טפסים נוקשים בעברית בלבד).",
            "  - ממשק קלט קולי מובנה ורציף: (מתחרים: לא קיים, דורש הקלדה ארוכה).",
            "  - אינטגרציה לניהול תורים ו-CRM: (מתחרים: טפסים סטטיים הדורשים חזרה טלפונית ידנית ממושכת).",
            "✔ תיקוף השערות במחקר:",
            "  - השערת דחיית בקשת הטלפון לסוף התהליך אומתה כמגבירה אמון; השערת ההקלטה הקולית אומתה כמפחיתה מאמץ הקלדה רפואית."
        ],
        is_rtl=True
    )

    # 5. Slide 5 (Part 5 - Business Model Canvas)
    create_content_slide(
        prs,
        "חלק 5: קנבס מודל עסקי - שלבים 1-4 בלבד",
        "קנבס מודל עסקי (חלק 5)",
        [
            ".1 פלחי לקוחות (Customer Segments):",
            "  - לקוחות ונפגעי תאונות עבודה ודרכים המחפשים בדיקת זכאות עצמאית ומהירה ותקשורת קלה מול עורך דין.",
            "  - משרדי עורכי דין לנזקי גוף המחפשים פניות מסוננות, מאומתות ומאורגנות מראש ללא שיחות סרק.",
            ".2 הצעת ערך (Value Propositions):",
            "  - ללקוח: בדיקה מהירה, אנונימית, ללא סיכון פיננסי, תזמון תור מיידי ושליטה מלאה בנתונים האישיים.",
            "  - לעורך הדין: קבלת תיק מסווג ומסונן היטב מראש, יומן פגישות מאורגן וחיסכון משמעותי בזמן הסינון.",
            ".3 ערוצים (Channels):",
            "  - פרסום ממומן ואורגני במנועי חיפוש (SEO/SEM), קבוצות פייסבוק מקצועיות וקהילות עובדים.",
            ".4 קשרי לקוחות (Customer Relationships):",
            "  - קשר מבוסס אמון ודיסקרטיות מלאה (אפשרות להשמדת נתונים בלחיצה), ליווי בצ'אט קשוב ללא לחץ משרדי."
        ],
        is_rtl=True
    )

    # 6. Slide 6 (Part 6 - Feature Spec & Prioritization)
    table_headers = ["שם הפיצ'ר", "ערך למשתמש", "עדיפות", "זמן (הערכה)", "מדד הצלחה"]
    table_rows = [
        ["צ'אט דינמי", "התאמת שאלות ומניעת סרבול ובלבול", "Must Have", "15 שעות", "אחוז השלמה מעל 85%"],
        ["זיהוי והקלטה קולית", "מניעת הקלדה רפואית קשה בנייד", "Must Have", "20 שעות", "שימוש קולי ב-50% מהתשובות"],
        ["מסך הסכמה וביטחון", "הסרת חשש משפטי ופיננסי ('0 ש\"ח')", "Must Have", "5 שעות", "ירידה בנטישה בשלב אישור הקשר"],
        ["לוח בקרה לעורך דין", "ארגון מובנה, תיאום תורים וסנכרון Outlook", "Must Have", "10 שעות", "זמן קריאת תיק פחות מ-3 דקות"],
        ["העלאת מסמכים חכמה", "העלאת אישורים וצילומים בקלות", "Should Have", "8 שעות", "העלאת מסמך ב-40% מהמקרים"]
    ]
    create_table_slide(
        prs,
        "חלק 6: אפיון ותעדוף פיצ'רים",
        "טבלת תעדוף ואפיון פיצ'רים (MoSCoW) (חלק 6)",
        table_headers,
        table_rows,
        is_rtl=True
    )

    # 7. Slide 7 (Part 7 & 8 - Functional Prototype & Usability)
    create_content_slide(
        prs,
        "חלק 7 וחלק 8: הפרוטוטייפ הפונקציונלי ודוח שמישות",
        "הפרוטוטייפ הפעיל ודוח שמישות ומדידה (חלק 7-8)",
        [
            "🔗 קישורים לפרוטוטייפ הפעיל:",
            "  - ממשק הלקוח וחיוג תורים: http://localhost:8000/client_intake.html",
            "  - לוח בקרה וניהול משרד עורך דין: http://localhost:8000/lawyer_dashboard.html",
            "📈 זרימת המשתמש בפרוטוטייפ (User Flow):",
            "  - בחירת שפה ומעבר לצ'אט -> תיאור קולי של התאונה -> בדיקת זכאות מיידית וזימון תור -> קבלת המידע בזמן אמת בלוח ה-CRM של עורך הדין וסנכרון מיידי לתיבת ה-Outlook.",
            "📊 דוח שמישות ומדידה (Usability Report):",
            "  - 2 סבבי בדיקות שמישות על 6 משתתפים שונים מקבוצת היעד ומחמיר אחד (עו\"ד).",
            "  - שיפורים שבוצעו: דחיית טלפון לסוף והוספת שסתום ביטחון | לוקליזציה של לחצני האישור | שדרוג הקלטה קולית רציפה.",
            " השוואת מדדים כמותיים (סבב 1 מול סבב 2):",
            "  - שיעור השלמת שאלון (Completion Rate): עלה מ-60% ל-90% בזכות הבהרת העלויות וההקלטה הקולית.",
            "  - זמן ממוצע להשלמה (Time to Complete): ירד מ-5 דקות ל-2 דקות בלבד הודות לממשק הקולי היעיל."
        ],
        is_rtl=True
    )

    # 8. Slide 8 (Personal Reflection / Experience)
    create_content_slide(
        prs,
        "עבודה אישית ותובנות מהפרויקט",
        "חלק 8: קשיים ותובנות אישיות בבניית המערכת",
        [
            "💡 קשיים טכנולוגיים ומוצריים שחוויתי במהלך הבנייה:",
            "  1. אינטגרציה דו-לשונית דינמית ומלאה:",
            "    * התמודדות עם ממשק דו-כיווני (RTL/LTR) במכשירים ניידים שונים, ותרגום דינמי של כפתורי הבחירה כמו \"OK\" / \"אישור\" / \"موافق\" בהתאם לשפת הלקוח כדי למנוע ממנו להתבלבל או להתייאש.",
            "  2. שיפור מנוע הקלטת הקול והמרת הדיבור לטקסט (Web Speech API):",
            "    * תכנון ופיתוח מנגנון הקלטה רציף שיודע לשרשר את הדיבור של הלקוח ולמנוע מחיקת טקסט רפואי קודם כאשר הלקוח לוקח הפסקה קלה בדיבור במהלך תיאור הפציעה.",
            "  3. פתרון מגבלות האבטחה באינטגרציית Microsoft Outlook:",
            "    * הדפדפנים המודרניים מונעים טעינה של מערכות מיקרוסופט ישירות בתוך Iframe עקב מגבלות אבטחה. פתרנו זאת על ידי פיתוח של מנגנון ייחודי בצד הלקוח הפותח חלונות נפרדים ומאובטחים העושים שימוש ב-SSO ובעוגיות הדפדפן הקיימות כדי לחבר את עורך הדין אוטומטית לחשבון שלו ללא צורך בהקשת סיסמה מחדש."
        ],
        is_rtl=True
    )

    prs.save("final_presentation_slides_he_updated.pptx")
    print("Hebrew presentation saved successfully.")

# ==========================================
# ARABIC PPTX BUILDER (8 SLIDES)
# ==========================================
def build_arabic_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Slide 1 (Title, Problem, Hypotheses)
    create_content_slide(
        prs,
        "جزء 1: وصف المشكلة والفرضيات",
        "وصف المشكلة والفرضيات الأساسية (جزء 1)",
        [
            "عنوان المشروع: جسر القدس القانوني (JLM Legal Bridge)",
            "عنوان فرعي: منصة رقمية ذكية لتسهيل وإدارة العلاقة بين المحامي والموكل، حجز المواعيد، وتقييم أهلية القضايا بالذكاء الاصطناعي.",
            "🔥 المشكلة الحقيقية (Pain حقيقي):",
            "  - عند المحامي: مكالمات ورسائل كثيرة غير مناسبة | تضييع وقت على زبائن غير مؤهلين | تكرار نفس الأسئلة.",
            "  - عند الزبون: لا يعرف إن كانت قضيته \"تستحق محامي\" | لغة قانونية معقدة | انتظار طويل وردود غير واضحة | خوف من نظرة المحامي وصعوبة التواصل بسبب المكان والزمان.",
            "🎯 الجمهور المستهدف واحتياجاته:",
            "  - الزبون: حاجته لمعرفة القانون بسهولة وسرعة عن بُعد دون إحراج (لذلك وفرنا الشات الذكي التفاعلي).",
            "  - المحامي: بحاجة لتسهيل وإدارة المواعيد وتواصل أسرع بملفات زبائن مصفاة وجاهزة.",
            "💡 الفرضيات الأولية (الفرضيات الأساسية):",
            "  - فرضية القيمة: فحص الأهلية الفوري والذكي عن بُعد سيزيد من استعداد الزبائن للتواصل بنسبة 40%.",
            "  - فرضية الاستخدام: استخدام التسجيل الصوتي يسهل وصف تفاصيل الحادث ويمنع هجر الاستبيان.",
            "  - فرضية الجدوى: دمج نظام الـ CRM والمواعيد سيوفر 60% من وقت الفرز والتنسيق الإداري للمكتب."
        ],
        is_rtl=True
    )

    # 2. Slide 2 (Part 2 - User Research & Insights)
    create_content_slide(
        prs,
        "جزء 2: ملخص بحث المستخدمين والرؤى المستخلصة",
        "نماذج المستخدمين والمهام وقصص المستخدم (جزء 2)",
        [
            "👥 نماذج المستخدمين (Personas):",
            "  - الزبون عادل (عامل بناء، 42 عاماً): يخشى التكاليف المسبقة ونظرة المحامي السلبية له، ويبحث عن تقييم فوري ومستقل.",
            "  - المحامي يوسف (شريك، 35 عاماً): يعاني من تضييع الوقت على زبائن غير مؤهلين وتكرار نفس الأسئلة.",
            "🔥 نقاط الألم المكتشفة (Pain Points):",
            "  - عند الزبون: ييأس ويبتعد عن المطالبة بحقه بسبب طول الإجراءات البيروقراطية والمجهود الكبير الذي يبذله لنجاح القضية.",
            "  - عند المحامي: هدر ساعات العمل الثمينة في مكالمات الفرز اليدوي وتكرار الأسئلة ذاتها.",
            "🎯 تحليل المهام المطلوبة (JTBD):",
            "  - المهمة 1 (الزبون): فحص أهلية قضيتي الطبية دون حرج أو التزام مالي (النجاح: إجابة واضحة \"قضيتك تستحق/لا تستحق\" خلال دقيقتين عن بعد).",
            "  - المهمة 2 (المحامي): تصفية واستلام ملفات موكلين مؤهلين وجاهزين بالوثائق (النجاح: استلام ملخص القضية الطبي بملف مدمج عبر لوحة التحكم).",
            "📝 قصص المستخدمين (User Stories):",
            "  - בתור זבון חרד שאינו בטוח אם המקרה שלו מוגדר כתאונת עבודה, כאשר אני מתבייש לפנות ישירות למשרד, אני רוצה להשתמש בצ'אט בוט אנונימי ונגיש בשפתי, כדי שאקבל אישור ראשוני לזכאותי ללא שיפוט.",
            "  - בתור משתמש רגיש המשתף פרטים רפואיים, כאשר אני מסיים את הבדיקה בצ'אט, אני רוצה לבחור באופן מפורש אם לשלוח את המידע או להשמיד אותו לצמיתות, כדי שארגיש שליטה מוחלטת במידע הפרטי שלי."
        ],
        is_rtl=True
    )

    # 3. Slide 3 (Part 3 - Market Research)
    create_content_slide(
        prs,
        "جزء 3: ملخص دراسة السوق والتوجهات",
        "الحلول الحالية، التوجهات التقنية وتحليل SWOT (جزء 3)",
        [
            "🔍 الحلول الحالية والفجوة في السوق:",
            "  - النماذج الرقمية الثابتة: طويلة ومعقدة، باللغة العبرية فقط، وتؤدي لهجر الموقع لغياب المنطق التفاعلي.",
            "  - التواصل المباشر والزيارة: تتطلب جهداً ووقتاً كبيراً، ومخاطرة مادية ونفسية تؤدي ليأس الزبون مبكراً.",
            "  - الفجوة المكتشفة: غياب منصة موحدة تجمع بين تقييم الأهلية التلقائي (AI) ولوحة تحكم المحامي لتنظيم تواصل أسرع دون تشتيت.",
            "📈 التوجهات التقنية الرائدة:",
            "  - الدردشة القانونية الذكية (Conversational Intake): أخذ تفاصيل الحوادث عبر دردشة ودية في الموبايل.",
            "  - تقييم الأهلية بالذكاء الاصطناعي (AI Eligibility): فحص إمكانية نجاح القضية تلقائياً بناءً على المعطيات الطبية.",
            "  - إدارة الأنظمة الموحدة (Unified Practice Management): ربط طلبات الزبائن مباشرة بجدول مواعيد وملفات المحامي (Outlook).",
            "📊 تحليل SWOT للمشروع:",
            "  - חוזקות (نقاط القوة): فحص الأهلية المخصص والذكي، توفير وقت الفرز للمحامي، تنظيم فوري للمواعيد وتكامل مع Outlook.",
            "  - חולשות (نقاط الضعف): اعتماد أداء المنصة على اتصال الإنترنت لدى هاتف الزبون.",
            "  - הזדמנויות (الفرص): التوسع مستقبلاً لدمج النظام مع WhatsApp Business API للوصول الأسهل للزبائن.",
            "  - איומים (التهديدات): مقاومة مكاتب المحاماة التقليدية للتطوير الرقمي أو استنساخ الفكرة من شركات أتمتة كبرى."
        ],
        is_rtl=True
    )

    # 4. Slide 4 (Part 4 - Product Vision & Value Prop)
    create_content_slide(
        prs,
        "جزء 4: رؤية المنتج وقيمة العرض الفريدة",
        "الرؤية، وقيمة العرض الفريدة وتثبيت الفرضيات (جزء 4)",
        [
            "🎯 رؤية المنتج: إزالة عوائق الخوف والتعقيد البيروقراطي، لتكون المنصة الرقمية الأكثر موثوقية وأماناً لربط المصابين بحقوقهم والتعويضات العادلة بأسرع وقت.",
            "💡 قيمة المنتج الفريدة (Unique Value Proposition):",
            "  - تقييم أهلية التعويض القانوني وتسهيل حجز المواعيد مع المحامي خلال دقيقتين فقط، بلغة بسيطة، مجاناً وبدون التزامات مسبقة – مع خصوصية كاملة وتحكم مطلق ببياناتك.",
            "⚔️ مقارنة المنافسين (النقاط التي نتفوق بها):",
            "  - توطين ودعم ثنائي اللغة ديناميكي بالكامل: (المنافسون: صفحات ثابتة وبالعبرية فقط).",
            "  - واجهة إدخال صوتي مدمجة وسلسة: (المنافسون: غير متوفرة، تعتمد على الكتابة الطويلة).",
            "  - تكامل نظام المواعيد والـ CRM للمحامي: (المنافسون: نماذج اتصال جامدة تتطلب اتصالاً هاتفياً لاحقاً).",
            "✔ التحقق من الفرضيات (Validation):",
            "  - فرضية تأخير طلب الهاتف للنهاية: أثبتت صحتها بمنع هجر الموقع وزيادة ثقة الزبون.",
            "  - فرضية التسجيل الصوتي: أثبتت نجاحها بتقليص جهد تعبئة المصطلحات الطبية للعمال."
        ],
        is_rtl=True
    )

    # 5. Slide 5 (Part 5 - Business Model Canvas)
    create_content_slide(
        prs,
        "جزء 5: نموذج العمل التجاري Canvas (المراحل 1-4 فقط)",
        "نموذج العمل التجاري (جزء 5)",
        [
            ".1 شرائح العملاء (Customer Segments):",
            "  - الزبائن والمصابون في حوادث العمل والطرق الباحثون عن تقييم فوري وتواصل سهل مع محامٍ.",
            "  - مكاتب المحاماة المتخصصة في قضايا التعويضات الباحثة عن ملفات عملاء مؤهلة ومصنفة تلقائياً.",
            ".2 القيمة المضافة (Value Propositions):",
            "  - للزبون: فحص سريع وسلس مجاناً، حجز مواعيد فوري، وحماية خصوصية تامة دون خطر مالي مسبق.",
            "  - للمحامي: ملف قضية متكامل ومصنف مع تقييم أولي، وجدول مواعيد منظم يمنع الفوضى والمكالمات المتكررة.",
            ".3 القنوات (Channels):",
            "  - التسويق الرقمي المستهدف (SEO/SEM) عبر محركات البحث، المجموعات المهنية وتجمعات العمال الرقمية.",
            ".4 علاقات العملاء (Customer Relationships):",
            "  - علاقة قائمة على الثقة والخصوصية المطلقة (إمكانية تدمير البيانات بلمسة واحدة)، وتسهيل حجز وتعديل المواعيد بمرونة."
        ],
        is_rtl=True
    )

    # 6. Slide 6 (Part 6 - Feature Spec & Prioritization)
    table_headers_ar = ["اسم الميزة", "القيمة المضافة للمستخدم", "الأولوية", "تقدير الوقت", "مقياس النجاح"]
    table_rows_ar = [
        ["الدردشة التفاعلية", "توجيه ذكي للأسئلة يمنع التشتت والتعقيد", "Must Have", "15 ساعة", "تجاوز نسبة إتمام الاستبيان لـ 85%"],
        ["التسجيل الصوتي", "إزالة عقبات الكتابة الطبية على الموبايل", "Must Have", "20 ساعة", "استخدام الصوت في 50% من الحالات"],
        ["شاشة الأمان المالي", "إزالة القلق المالي ('0 شيكل مسبقاً')", "Must Have", "5 ساعات", "انخفاض هجر الاستبيان في خطوة الاتصال"],
        ["لوحة تحكم المحامي", "تنظيم الملفات وتنسيق المواعيد وتكامل Outlook", "Must Have", "10 ساعات", "تصفح وقراءة ملف القضية بأقل من 3 دقائق"],
        ["رفع الملفات الذكي", "رفع التقارير الطبية والمستندات بمرونة", "Should Have", "8 ساعات", "رفع مستندات طبية في 40% من الحالات"]
    ]
    create_table_slide(
        prs,
        "جزء 6: مواصفات وتحديد أولويات الميزات (MVP)",
        "جدول تحديد وتفضيل الميزات الأساسية حسب نموذج MoSCoW (جزء 6)",
        table_headers_ar,
        table_rows_ar,
        is_rtl=True
    )

    # 7. Slide 7 (Part 7 & 8 - Functional Prototype & Usability)
    create_content_slide(
        prs,
        "جزء 7 وجزء 8: النموذج الأولي العملي وتقرير سهولة الاستخدام",
        "النموذج الأولي العملي وتقرير سهولة الاستخدام والقياس (جزء 7-8)",
        [
            "🔗 روابط تجربة النموذج الأولي الفعال:",
            "  - الاستبيان وحجز المواعيد للزبون: http://localhost:8000/client_intake.html",
            "  - لوحة تحكم المحامي وإدارة المكاتب: http://localhost:8000/lawyer_dashboard.html",
            "📈 خطوات تدفق الموكل والمحامي في النظام:",
            "  - اختيار اللغة والبدء بالدردشة -> التسجيل الصوتي لشرح الإصابة -> الحصول على تقييم أولي وحجز موعد فوري -> استلام المحامي للملف في لوحة التحكم وتنسيق الموعد تلقائياً مع Outlook.",
            "📊 تقرير سهولة الاستخدام والقياس بالأرقام (Usability Report):",
            "  - جولتا اختبارات شמישות على 6 مشاركين من الجمهور المستهدف ومحامٍ واحد.",
            "  - التحسينات المنجزة: تأخير طلب الهاتف للنهاية ووضع رسالة أمان ('0 شيكل') | تعريب زر التأكيد حسب الواجهة | تفعيل التصفية التلقائية.",
            " مقارنة المقاييس الرقمية (الجولة 1 مقابل الجولة 2):",
            "  - معدل إكمال الاستبيان (Completion Rate): ارتفع من 60% إلى 90% بفضل رسائل الأمان والصوت.",
            "  - متوسط وقت الإكمال (Time to Complete): انخفض من 5 دقائق إلى دقيقتين فقط بفضل واجهة الإدخال الصوتي وسرعة حجز المواعيد."
        ],
        is_rtl=True
    )

    # 8. Slide 8 (Personal Reflection / Experience)
    create_content_slide(
        prs,
        "التجربة الشخصية والتحديات التي واجهتها",
        "العمل الشخصي والتحديات والصعوبات التقنية التي واجهتها أثناء البناء (شريحة 8)",
        [
            "💡 الصعوبات التقنية والمنتجية التي واجهتها أثناء البناء:",
            "  1. الدمج ثنائي اللغة الديناميكي والكامل:",
            "    * الصعوبة في الحفاظ على تجربة مستخدم موحدة للدردشة والمستندات، وتجاوب أزرار الاستبيان ديناميكياً عند تحويل اللغة لتجنب تشتيت الزبون.",
            "  2. تحسين تقنية تحويل الصوت إلى نصوص (Web Speech API):",
            "    * تطلب الأمر ضبط إعدادات التسجيل والتعرف الصوتي لربط وتثبيت النصوص الطويلة بشكل متواصل لمنع مسح التفاصيل السابقة عندما يتوقف الزبون لالتقاط أنفاسه.",
            "  3. دمج وتأمين بوابة خدمات Microsoft Outlook للمحامي:",
            "    * واجهنا عائق منع التضمين (X-Frame-Options) للبريد الإلكتروني الحقيقي بداخل الموقع. قمنا بحل ذلك عن طريق بناء شريط أدوات Microsoft 365 الجانبي الذي يفتح التطبيقات بنوافذ مستقلة وآمنة لتسجيل الدخول التلقائي (SSO) عبر حساب ويندوز النشط للمحامي."
        ],
        is_rtl=True
    )

    prs.save("final_presentation_slides_ar_updated.pptx")
    print("Arabic presentation saved successfully.")

if __name__ == '__main__':
    build_hebrew_presentation()
    build_arabic_presentation()
